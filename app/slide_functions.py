import os
import json

from enum import Enum
from pptx.util import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE


# Default Microsoft Office Slide Layouts
class DefaultSlideLayouts(Enum):
    TITLE_SLIDE = 0
    TITLE_AND_CONTENT_SLIDE = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    COMPARISON = 4
    TITLE_ONLY = 5
    BLANK = 6
    CONTENT_WITH_CAPTION = 7
    PICTURE_WITH_CAPTION = 8
    TITLE_AND_VERTICAL_TEXT = 9
    VERTICAL_TITLE_AND_TEXT = 10


################################
#    Slide creation functions    #
################################

def create_title_slide(slide_content, pres):
    """
    Function that create a title slide given some slide content
    Example of slide_content:
    ['level_0',
        'title_0',
        'My super title',
        'plain_text',
        'My subtitle'
    ]
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['TITLE_SLIDE'].value]
    slide = pres.slides.add_slide(slide_layout)

    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be a subtitle
        elif ("plain_text" in content):
            try:
                # print(content, slide_content[i+1])
                slide.placeholders[1].text = "\n".join(slide_content[i + 1:])
            except IndexError:
                print('The text is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no such placeholder.")
                print(f"Could not write {content} into the presentation")

    return pres


def create_section_header_slide(slide_content, pres):
    """
    Function that create a section header slide given some slide content
    Example of slide_content:
    ['level_1',
        'title_1',
        'Part 1 of presentation',
        'plain_text',
        'Small text'
    ]
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['SECTION_HEADER'].value]
    slide = pres.slides.add_slide(slide_layout)
    plain_text = []
    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be subheader
        elif ("plain_text" in content):
            try:
                # print(content, slide_content[i+1])
                plain_text.append(slide_content[i + 1])
            except IndexError:
                print('The text is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
    try:
        text = "\n".join("".join(plain_text).split('. '))
        slide.placeholders[1].text = text
    except:
        print("This slide has no such placeholder.")
        print(f"Could not write {text} into the presentation")
    return pres


def create_title_and_content_slide(slide_content, pres):
    """
    Function that create a title and content slide given some slide content
    Example of slide_content:
    ['level_n',
        'title_n',
        'My super title',
        'plain_text',
        'My first text',
        'Some other piece of info',
        'Meaningful sentence'
    ]
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['TITLE_AND_CONTENT_SLIDE'].value]
    slide = pres.slides.add_slide(slide_layout)
    plain_text = []
    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be the content (usually plain text)
        elif ("plain_text" in content):
            try:
                # print(content, slide_content[i+1])
                plain_text.append(slide_content[i + 1])
            except IndexError:
                print('The text is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
    try:
        text = "\n".join("".join(plain_text).split('. '))
        slide.placeholders[1].text = text
    except:
        print("This slide has no such placeholder.")
        print(f"Could not write {text} into the presentation")

    return pres


def create_two_content_slide(slide_content, pres):
    """
    Function that create a two content slide given some slide content
    Example of slide_content:
    ['level_n',
        'title_n',
        'My super title',
        'plain_text_1',
        'My first text',
        'My second text',
        'plain_text_2',
        'My third text'
    ]
    """
    visited = False

    slide_layout = pres.slide_layouts[DefaultSlideLayouts['TWO_CONTENT'].value]
    slide = pres.slides.add_slide(slide_layout)

    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be plain text (here two contents)
        elif ("plain_text" in content):
            if not (visited):
                visited = True
                # first content
                try:
                    end_of_text = slide_content.index("plain_text_2")
                    # print(content, slide_content[i+1:end_of_text])
                    slide.placeholders[1].text = "\n".join(slide_content[i + 1:end_of_text])
                except IndexError:
                    print('The text is empty. Check the slide content.')
                    for content in slide_content:
                        print(content)
                except:
                    print("This slide has no such placeholder.")
                    print(f"Could not write {content} into the presentation")
            else:
                # second content
                try:
                    # print(content, slide_content[i+1:])
                    slide.placeholders[2].text = "\n".join(slide_content[i + 1:])
                except IndexError:
                    print('The text is empty. Check the slide content.')
                    for content in slide_content:
                        print(content)
                except:
                    print("This slide has no such placeholder.")
                    print(f"Could not write {content} into the presentation")

    return pres


def create_comparison_slide(slide_content, pres):
    """
    Function that create a two content slide given some slide content
    Example of slide_content:
    ['level_n',
        'title_n',
        'My super title',
        'header_1',
        'First subtitle',
        'plain_text_1',
        'My first text',
        'My second text',
        'header_2',
        'Second subtitle',
        'plain_text_2',
        'My third text',
        'My fourth text'
    ]
    """
    visited = False

    slide_layout = pres.slide_layouts[DefaultSlideLayouts['COMPARISON'].value]
    slide = pres.slides.add_slide(slide_layout)

    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # the two subtitles
        elif ("header" in content):
            if (visited):
                visited = False
                # second subtitle
                try:
                    # print(content, slide_content[i+1])
                    slide.placeholders[3].text = slide_content[i + 1]
                except IndexError:
                    print('The text is empty. Check the slide content.')
                    for content in slide_content:
                        print(content)
                except:
                    print("This slide has no such placeholder.")
                    print(f"Could not write {content} into the presentation")
            else:
                visited = True
                # first subtitle
                try:
                    # print(content, slide_content[i+1])
                    slide.placeholders[1].text = slide_content[i + 1]
                except IndexError:
                    print('The text is empty. Check the slide content.')
                    for content in slide_content:
                        print(content)
                except:
                    print("This slide has no such placeholder.")
                    print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be plain text (here two contents)
        elif ("plain_text" in content):
            if (visited):
                # first content
                try:
                    end_of_text = slide_content.index("header_2")
                    # print(content, slide_content[i+1:end_of_text])
                    slide.placeholders[2].text = "\n".join(slide_content[i + 1:end_of_text])
                except IndexError:
                    print('The text is empty. Check the slide content.')
                    for content in slide_content:
                        print(content)
                except:
                    print("This slide has no such placeholder.")
                    print(f"Could not write {content} into the presentation")
            else:
                # second content
                try:
                    # print(content, slide_content[i+1:])
                    slide.placeholders[4].text = "\n".join(slide_content[i + 1:])
                except IndexError:
                    print('The text is empty. Check the slide content.')
                    for content in slide_content:
                        print(content)
                except:
                    print("This slide has no such placeholder.")
                    print(f"Could not write {content} into the presentation")

    return pres


def create_title_only_slide(slide_content, pres):
    """
    Function that create a title only slide given some slide content
    Example of slide_content:
    ['level_01',
        'title_01',
        'My super title',
    ]
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['TITLE_ONLY'].value]
    slide = pres.slides.add_slide(slide_layout)

    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

    return pres


def create_content_with_caption_slide(slide_content, pres, imgspath):
    """
    Function that create a content with caption slide given some slide content
    Example of slide_content:
    ['level_n',
        'title_n',
        'My super title',
        'content_item',
        'My content can be anything'
        'plain_text',
        'Caption text'
    ]
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['CONTENT_WITH_CAPTION'].value]
    slide = pres.slides.add_slide(slide_layout)
    plain_text = []
    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # set the content
        elif ("content_item" in content):
            # picture or plain text only for now
            # charts, tables and stuff coming later
            try:
                # print(content, slide_content[i+1])
                picture = slide.placeholders[1].insert_picture(slide_content[i + 1])
            except IndexError:
                print('The picture is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except FileNotFoundError:
                print('File not found. Placed a default picture instead.')
                picture = slide.placeholders[1].insert_picture(imgspath + 'default_picture.jpg')
            except:
                print("This slide has no such placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be plain text
        elif ("plain_text" in content):
            try:
                # print(content, slide_content[i+1])
                plain_text.append(slide_content[i + 1])
            except IndexError:
                print('The text is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
    try:
        text = "\n".join("".join(plain_text).split('. '))
        slide.placeholders[1].text = text
    except:
        print("This slide has no such placeholder.")
        print(f"Could not write {text} into the presentation")

    return pres


def create_picture_with_caption_slide(slide_content, pres, imgspath):
    """
    Function that create a picture with caption slide given some slide content
    Example of slide_content:
    ['level_n',
        'title_n',
        'My super title',
        'picture_item',
        'picture path'
        'plain_text',
        'Caption text'
    ]
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['PICTURE_WITH_CAPTION'].value]
    slide = pres.slides.add_slide(slide_layout)
    plain_text = []
    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # set the picture
        elif ("picture_item" in content):
            try:
                # print(content, slide_content[i+1])
                picture = slide.placeholders[1].insert_picture(slide_content[i + 1])
            except IndexError:
                print('The picture is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except FileNotFoundError:
                print('File not found. Placed a default picture instead.')
                picture = slide.placeholders[1].insert_picture(imgspath + 'default_picture.jpg')
            except:
                print("This slide has no such placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be plain text
        elif ("plain_text" in content):
            try:
                # print(content, slide_content[i+1])
                plain_text.append(slide_content[i + 1])
            except IndexError:
                print('The text is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
    try:
        text = "\n".join("".join(plain_text).split('. '))
        slide.placeholders[1].text = text
    except:
        print("This slide has no such placeholder.")
        print(f"Could not write {text} into the presentation")

    return pres


def create_title_and_vertical_text_slide(slide_content, pres):
    """
    Same function as title_and_content but with vertical layout
    TO DO: add a keyword to say the content should be vertical
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['TITLE_AND_VERTICAL_TEXT'].value]
    slide = pres.slides.add_slide(slide_layout)
    plain_text = []
    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be plain text
        elif ("plain_text" in content):
            try:
                # print(content, slide_content[i+1])
                plain_text.append(slide_content[i + 1])
            except IndexError:
                print('The text is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
    try:
        text = "\n".join("".join(plain_text).split('. '))
        slide.placeholders[1].text = text
    except:
        print("This slide has no such placeholder.")
        print(f"Could not write {text} into the presentation")

    return pres


def create_vertical_title_and_text_slide(slide_content, pres):
    """
    Same function as title_and_content but with vertical layout
    TO DO: add a keyword to say the content should be vertical
    """
    slide_layout = pres.slide_layouts[DefaultSlideLayouts['VERTICAL_TITLE_AND_TEXT'].value]
    slide = pres.slides.add_slide(slide_layout)
    plain_text = []
    for i, content in enumerate(slide_content):
        # set the title
        if ('title_' in content):
            try:
                # print(content, slide_content[i+1])
                title = slide.shapes.title
                title.text = slide_content[i + 1]
            except IndexError:
                print('The title is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
            except:
                print("This slide has no title placeholder.")
                print(f"Could not write {content} into the presentation")

        # logically, the rest of the slide should be plain text
        elif ("plain_text" in content):
            try:
                # print(content, slide_content[i+1])
                plain_text.append(slide_content[i + 1])
            except IndexError:
                print('The text is empty. Check the slide content.')
                for content in slide_content:
                    print(content)
    try:
        text = "\n".join("".join(plain_text).split('. '))
        slide.placeholders[1].text = text
    except:
        print("This slide has no such placeholder.")
        print(f"Could not write {text} into the presentation")

    return pres


def detect_layout_and_create_slide(slide_contents, pres, imgspath):
    """
    Function that automatically detect the slide layout to use given its content.
    Parameters:
        - slide_contents: list of slide contents
        - pres: the presentation to add slides into
    """
    for slide_content in slide_contents:
        if (int(slide_content[0].split("_")[1]) == 0):
            print("Creating a title slide\n")
            create_title_slide(slide_content, pres)
        elif (int(slide_content[0].split("_")[1]) == 1):
            print("Creating a section header slide\n")
            create_section_header_slide(slide_content, pres)
        elif (int(slide_content[0].split("_")[1]) >= 2):
            try:
                idx = slide_content.index('plain_text')
                if (len(slide_content) - idx == 3):
                    if ("headers" in slide_content):
                        print("Creating a comparison slide\n")
                        create_comparison_slide(slide_content, pres)
                    else:
                        print("Creating a two content slide\n")
                        create_two_content_slide(slide_content, pres)
            except ValueError:
                pass
            if ("picture_item" in slide_content):
                print("Creating a picture with caption slide\n")
                create_picture_with_caption_slide(slide_content, pres, imgspath)
            elif ("content" in slide_content):
                print("Creating a content with caption slide\n")
                create_content_with_caption_slide(slide_content, pres)
            else:
                print("Creating a title and content slide\n")
                create_title_and_content_slide(slide_content, pres)
        else:
            print("Slide layout unknown")

    return pres